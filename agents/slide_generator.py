"""Slide Generator — builds slide outline and exports a .pptx."""
import json, re, io
from langchain_core.prompts import ChatPromptTemplate
from pptx import Presentation
from pptx.util import Inches, Pt
from .llm_factory import get_llm

PROMPT = ChatPromptTemplate.from_template(
    """You are a slide-deck designer for an academic talk.

Using the summary and extracted concepts below, produce a JSON slide deck.

Return STRICT JSON:
{{
  "title": "Deck title",
  "subtitle": "1-line subtitle",
  "slides": [
    {{ "title": "Slide Title", "bullets": ["bullet 1", "bullet 2", "bullet 3"] }}
  ]
}}

Requirements:
- 8 to 10 content slides total (plus the title slide, which is separate).
- Cover: Motivation, Problem, Related Work, Method/Approach, Architecture/Algorithm,
  Experiments & Datasets, Results, Limitations, Conclusion, Future Work.
- 3-5 concise bullets per slide. Max ~15 words per bullet.
- Be specific to THIS paper.

Summary:
{summary}

Concepts:
{concepts}

Return only JSON.
"""
)


def _parse_json(text: str) -> dict:
    text = re.sub(r"^```(?:json)?|```$", "", text.strip(), flags=re.MULTILINE).strip()
    try:
        return json.loads(text)
    except Exception:
        m = re.search(r"\{.*\}", text, flags=re.DOTALL)
        return json.loads(m.group(0))


def generate_slide_plan(summary: dict, concepts: dict) -> dict:
    llm = get_llm(temperature=0.4)
    chain = PROMPT | llm
    resp = chain.invoke({
        "summary": json.dumps(summary, indent=2),
        "concepts": json.dumps(concepts, indent=2),
    })
    return _parse_json(resp.content)


def build_pptx(plan: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)
    s.shapes.title.text = plan.get("title", "Research Paper")
    if len(s.placeholders) > 1:
        s.placeholders[1].text = plan.get("subtitle", "")

    # Content slides
    for sl in plan.get("slides", []):
        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = sl.get("title", "")
        body = s.placeholders[1].text_frame
        body.clear()
        bullets = sl.get("bullets", [])
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.font.size = Pt(20)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
